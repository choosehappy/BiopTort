import os
from typing import Tuple, Optional
from skimage.morphology import remove_small_holes, remove_small_objects, binary_dilation, disk
from skimage.measure import regionprops, find_contours
from skimage.segmentation import active_contour
from skimage.filters import rank
from skimage import color
from shapely import LineString, Polygon
from scipy.ndimage import binary_dilation, distance_transform_edt, binary_fill_holes, gaussian_filter1d
from scipy.interpolate import interp1d
import openslide
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
import glob
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
import PIL
from tqdm import tqdm
import cv2
import csv
import argparse
from bioptort.mask_utils import split_binary_cnb_mask

# ------------------- CONSTANTS -------------------
WORKING_MAGNIFICATION = 0.625

# ------------------- UTILS -------------------
def get_base_magnification(osh: openslide.OpenSlide) -> float:
    val = osh.properties.get("openslide.objective-power") \
          or osh.properties.get("aperio.AppMag")
    
    if val is None:
        raise ValueError("Base magnification not found in slide properties. Please provide it manually using the --override_base_magnification argument.")
    
    return float(val)

def get_downsample_factor(base_magnification: float) -> float:
    return base_magnification / WORKING_MAGNIFICATION

def get_best_level_for_downsample(osh: openslide.OpenSlide, downsample_factor: float) -> Tuple[int, bool]: 
        relative_down_factors_idx=[np.isclose(i/downsample_factor,1,atol=.01) for i in osh.level_downsamples]
        level=np.where(relative_down_factors_idx)[0]
        if level.size:
            return (level[0], True)
        else:
            return (osh.get_best_level_for_downsample(downsample_factor), False)

def resize_image(img: Image.Image, basewidth=None, scalefactor=0.5) -> Image.Image:
    """Resize an image to a specified width while maintaining aspect ratio.

    Args:
        img (Image): the image to resize
        basewidth (int, optional): the width to resize to. Defaults to None.
        scalefactor (float, optional): the scale factor to resize by. Defaults to 0.5. Only applied if basewidth is not specified.

    Returns:
        Image: the resized image
    """
    if basewidth is None:
        return img.resize((int(img.size[0] * scalefactor), int(img.size[1] * scalefactor)))

    wpercent = (basewidth/float(img.size[0]))
    hsize = int((float(img.size[1])*float(wpercent)))
    
    return img.resize((basewidth, hsize))

def get_downsampled_image(osh: openslide.OpenSlide, override_base_magnification: float=None) -> Image.Image:
    base_mag = override_base_magnification or get_base_magnification(osh)

    target_downsample_factor = get_downsample_factor(base_mag)
    level, isexactlevel = get_best_level_for_downsample(osh, target_downsample_factor)

    region = osh.read_region((0, 0), level, osh.level_dimensions[level])

    if isexactlevel:
        return resize_image(region, scalefactor=1.0)
    else:
        current_dims = osh.level_dimensions[level]
        target_dims = tuple(np.rint(np.asarray(osh.level_dimensions[0]) / target_downsample_factor).astype(int))
        scalefactor = target_dims[0] / current_dims[0]
        resized_region = resize_image(region, scalefactor=scalefactor)
        return resized_region


# ------------------- LOOP REMOVAL -------------------
def segments_intersect(s1, s2):
    return LineString(s1).intersects(LineString(s2))

def get_segments_from_points(points):
    segments = []
    for i in range(len(points)-1):
        segments.append(np.array([points[i], points[i+1]]))
    return segments

def get_points_from_segments(segments):
    points = []
    for s in segments:
        points.append(s[0])
    points.append(segments[-1][1])
    return np.array(points)

def remove_loops(points):
    stack = []
    loops = []

    points_copy = np.copy(points)

    segments = get_segments_from_points(points)
    for i, s1 in enumerate(segments):
        for j, s2 in enumerate(segments):
            if segments_intersect(s1, s2): 
                if j > i + 1:   # if the segments are not adjacent we add the first segment to the stack
                    stack.append(i)
                elif j < i - 1:
                    loopstart = stack.pop()
                
                    if len(stack) == 0:
                        loops.append((loopstart, i))

    
    for i in reversed(loops):
        points_copy = np.delete(points_copy, slice(i[0]+1, i[1]+1), axis=0) # not very efficient but it works

    return points_copy

# ---------------------------------------------------
def remove_background_points(points, mask):
    try:
        return np.array([p for p in points if mask[int(p[0]), int(p[1])]])
    except IndexError:
        raise IndexError("L contains points that exceed the boundaries of the image. Your image likely contains artifacts near the edges.")

def apply_gaussian_filter(point_array, sigma=1):
    X_values = point_array[:, 0]
    Y_values = point_array[:, 1]
    X_values = gaussian_filter1d(X_values, sigma=sigma)
    Y_values = gaussian_filter1d(Y_values, sigma=sigma)

    # --- genepy denoising method
    # noisy_curve = curves.Curve((X_values, Y_values))
    # denoised_curve = noisy_curve.denoise()
    # return denoised_curve.coors[:, :2]

    return np.vstack((X_values, Y_values)).T

def interpolate_points(points, num_points):
    # Linear length along the line:
    distance = np.cumsum( np.sqrt(np.sum( np.diff(points, axis=0)**2, axis=1 )) )
    distance = np.insert(distance, 0, 0)/distance[-1]

    alpha = np.linspace(distance.min(), int(distance.max()), num_points * 10)
    interpolator = interp1d(distance, points, kind='slinear', axis=0)
    interpolated_points = interpolator(alpha)

    out_x = interpolated_points.T[0]
    out_y = interpolated_points.T[1]

    return np.vstack((out_x, out_y)).T

def addimagetoslide(slide, image_stream, left, top, height, width, resize=1.0, comment=''):
    # res = cv2.resize(img , None, fx=resize,fy=resize ,interpolation=cv2.INTER_CUBIC) #since the images are going to be small, we can resize them to prevent the final pptx file from being large for no reason
    # image_stream = BytesIO()
    # Image.fromarray(res).save(image_stream,format="PNG")
    slide.shapes.add_picture(image_stream, left, top, width, height)
    txBox = slide.shapes.add_textbox(left, Inches(1), width, height)
    tf = txBox.text_frame
    tf.text = comment

def move_slide(slides, slide, new_idx):
    slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])

def closest_farthest_points(origin, point_set):
    """Calculate the closest and furthest point in a set from a specified origin.

    Args:
        origin (list): list of length 2 
        point_set (list): list of shape (N,2)

    Returns:
        tuple: a tuple containing the closest point and the farthest point.
    """    
    distance_metric = lambda x: (x[0] - origin[0])**2 + (x[1] - origin[1])**2 # no need to sqrt for comparing distances
    closest = min(point_set, key=distance_metric)
    farthest = max(point_set, key=distance_metric)

    return closest, farthest 

def filter_contours(contours, bounds, angle, L):
    """Filter out contours that are not touching tissue.

    Args:
        contours (ndarray): the original set of contours
        bounds (tuple): the two points marking the 10th and 90th percentile of the snake.
        angle (float): the angle of orientation of the tissue.
        L (ndarray): the snake.
    """
    filtered_contours = []
    for contour in contours:
        if LineString(L).intersects(Polygon(contour)):
            for point in contour:
                within_bound1 = np.sign(angle) * (np.sin(angle) * (point[0] - bounds[0][0]) + np.cos(angle) * (point[1] - bounds[0][1])) > 0
                within_bound2 = np.sign(angle) * (np.sin(angle) * (point[0] - bounds[1][0]) + np.cos(angle) * (point[1] - bounds[1][1])) < 0

                if within_bound1 and within_bound2:
                    filtered_contours.append(contour)
                    break
            
    return filtered_contours

def adjust_starting_line(b, l, angle):
    """Adjust the location of each point until it is touching tissue. 

    Args:
        b (ndarray): the binary mask
        l (ndarray): the coordinates of L 
        angle (float): the angle of orientation of the tissue.

    Returns:
        ndarray: the new coordinates of the starting line.
    """    
    new_l = []

    # calculate the movement vector from the angle in row col space
    movement_vect = np.array([np.sin(-1 * angle), np.cos(angle)])

    def check_bounds(point, b):   
        return point[0] >= 0 and point[1] >= 0 and point[0] < b.shape[0] and point[1] < b.shape[1]
    
    for point in l:
        hit = False
        n = np.copy(point)
        p = np.copy(point)

        # if the point is already touching tissue, add it to the new list.
        while not hit and (check_bounds(n, b) and check_bounds(p, b)):    # check if the point is within the bounds of the image
            if b[int(n[0]), int(n[1])]:
                new_l.append(n)
                hit = True
            elif b[int(p[0]), int(p[1])]:
                new_l.append(p)
                hit = True
            else:
                n -= movement_vect * 3
                p += movement_vect * 3

        # if not hit:
        #     new_l.append(point)
            
    return np.array(new_l)

def calculate_grade(xy_tortuosity, num_skips):
    """Calculate the grade of the tortuosity based on the number of skips and the tortuosity.

    Args:
        xy_tortuosity (float): the tortuosity of the snake, [1, infinity)
        num_skips (int): the number of skips detected

    Returns:
        int: the predicted grade
    """
    # thresh = 1.3
    thresh = 1.076

    if num_skips == 0:
        if xy_tortuosity < thresh:
            return 1
        else:
            return 2
    else:
        if xy_tortuosity < thresh:
            return 3
        else:
            return 4

def resize_image(img: Image.Image, basewidth=None, scalefactor=0.5) -> Image.Image:
    """Resize an image to a specified width while maintaining aspect ratio.

    Args:
        img (Image): the image to resize
        basewidth (int, optional): the width to resize to. Defaults to None.
        scalefactor (float, optional): the scale factor to resize by. Defaults to 0.5. Only applied if basewidth is not specified.

    Returns:
        Image: the resized image
    """
    if basewidth is None:
        return img.resize((int(img.size[0] * scalefactor), int(img.size[1] * scalefactor)))

    wpercent = (basewidth/float(img.size[0]))
    hsize = int((float(img.size[1])*float(wpercent)))
    
    return img.resize((basewidth, hsize))

def generate_mask(img: np.ndarray, disk_size=2, threshold=210) -> np.ndarray:
    """
    Generates a binary mask from an input image by applying a series of image processing steps. Inspired by https://github.com/choosehappy/HistoQC/blob/13a9439a482595fa45469837b05d2d0e326843a3/histoqc/LightDarkModule.py#L146

    Parameters:
    img (Image): The input image from which the mask is to be generated.
    disk_size (int, optional): The size of the disk-shaped structuring element used for the minimum filter. Default is 10000.
    threshold (int, optional): The threshold value used to create the binary mask. Default is 200.

    Returns:
    numpy.ndarray: A binary mask where pixels below the threshold are set to True, and others are set to False.
    """
    # naive approach: threshold the red channel of the image
    # return np.asarray(img)[:, :, 0] < 200
    img = color.rgb2gray(np.array(img))
    img = (img * 255).astype(np.uint8)
    selem = disk(disk_size)
    imgfilt = rank.minimum(img, selem)

    return imgfilt < threshold  # return the binary mask

def process_mask(mask: PIL.Image) -> np.ndarray:
    b = remove_small_holes(mask, 1000)
    b = remove_small_objects(b, 100)

    return binary_fill_holes(
        binary_dilation(b, iterations=1))


def compute_snake(b: np.ndarray):
    # --- Hyperparameters ---
    pixel_lengths_between_points = 6    # the number of pixel-lengths between each point on the initial line.
    # --- active contour parameters ---
    alpha = 0.001
    beta = 0.001 # 0.4
    gamma = 0.001
    w_line = 50
    max_px_move = 1
    w_edge = 0
    max_num_iter = 1000
    convergence = .01
    # -----------------------------

    # basewidth = 1000
    # wpercent = (basewidth/float(img.size[0]))
    # hsize = int((float(img.size[1])*float(wpercent)))
    # img = img.resize((basewidth, hsize))

    # find contours
    contours = find_contours(b)
    

    # compute the distance transform and normalize to the range [0, 0.5]
    eb = distance_transform_edt(b)
    eb = cv2.normalize(eb, None, 0, 1, cv2.NORM_MINMAX) / 2

    # perform a similar computation on the inverted mask. Negate and normalize to the same range.
    e_inv_b = distance_transform_edt(np.logical_not(b))
    e_inv_b = (1 - cv2.normalize(e_inv_b, None, 0, 1, cv2.NORM_MINMAX)) / 2

    # sum the two distributions
    eb_total = eb + e_inv_b

    props = regionprops(b.astype(np.uint8))

    if len(props) == 0:
        raise ValueError("No tissue regions found in the binary mask.")

    rp = props[0]

    print(rp.orientation)
    if rp.orientation <= 0:
        rc_start = (rp.bbox[0], rp.bbox[3])
    else:
        rc_start = (rp.bbox[0], rp.bbox[1])

    contour_points = []
    for contour in contours:
        contour_points.extend(contour)

    closest, farthest = closest_farthest_points(rc_start, contour_points)

    dist = np.linalg.norm(closest - farthest)
    num_points = int(dist / pixel_lengths_between_points)
    r = np.linspace(closest[0], farthest[0], num_points)
    c = np.linspace(closest[1], farthest[1], num_points)
    
    init = np.array([r, c]).T
    init_adjusted = adjust_starting_line(b, init, rp.orientation)

    print(f'Region Props bounding box: {rp.bbox}')
    snake = active_contour(eb_total, init_adjusted, alpha=alpha, beta=beta, gamma=gamma, w_line=w_line, max_px_move=max_px_move,
                    w_edge=w_edge, max_num_iter=max_num_iter, convergence=convergence, boundary_condition='fixed')
    
    return snake, contours, rp, b, num_points
    

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--im_path', type=str, required=True, help="The directory containing (openslide-compatible) CNB WSIs.")
    parser.add_argument('--ppt_save', type=str, required=True, help="The path to save the powerpoint presentation, including the filename.")
    parser.add_argument('--csv_save', type=str, required=False, help="Optional path to save a csv file containing quantitative results, including the filename.")
    parser.add_argument('--mask_path', type=str, default=None, required=False, help="Optional path to existing png tissue masks for the CNB WSIs with matching filenames.")
    parser.add_argument('--sort', '-s', default=False, action='store_true', help="Sort the slides in order of highest to lowest tortuosity.")
    parser.add_argument('--gauss_sigma', type=float, default=20, help="The sigma value for the gaussian filter.")
    parser.add_argument('--override_base_magnification', type=float, default=20.0, help="Manually specify the base magnification of the WSI if it cannot be read from the file properties.")
    parser.add_argument('--multiple_sections', '-m', default=False, action='store_true', help="Experimental: Enable handling of multiple sections per slide. May not work as intended.")
    args = parser.parse_args()
    files = sorted(glob.glob(os.path.join(args.im_path, "*.svs")))

    if args.mask_path:
        mask_filepaths = sorted(glob.glob(os.path.join(args.mask_path, "*.png")))
    # ------------- Main script -------------
    # init presentation
    ppt = Presentation()
    tort_values = []
    sorted_fns = []
    sorted_skips = []
    sorted_pred_grades = []
    sorted_core_mask_indices = []

    for fn_i, fn in tqdm(enumerate(files)):
        # ------------- LOAD IMAGE -------------
        osh = openslide.OpenSlide(fn)
        print(f'Working on file: {fn}')
        print(f'Level dimensions: {osh.level_dimensions}')
        image = get_downsampled_image(osh, override_base_magnification=args.override_base_magnification)

        # Load the mask if it exists
        if args.mask_path:
            mask = Image.open(mask_filepaths[fn_i])
            mask = resize_image(mask, basewidth=image.size[0])
            binary_mask = np.asarray(mask)
            binary_mask = np.where(binary_mask > 0, 1, 0)
        else:
            npy_img = np.array(image)[:,:,:3]
            binary_mask = generate_mask(npy_img)

        binary_mask = process_mask(binary_mask)
        core_masks = split_binary_cnb_mask(binary_mask) if args.multiple_sections else [binary_mask]

        for core_mask_i, core_mask in enumerate(core_masks):

            # ------------- COMPUTE SNAKE -------------
            # ... and other image processing
            try:
                snake, contours, rp, b, num_points = compute_snake(core_mask)
            except ValueError as e:
                print(f"Skipping core {core_mask_i} in file {fn} due to error: {e}")
                continue


            # ------------- REMOVE LOOPS -------------
            snake = remove_loops(snake)
            # ------------- SNAKE SMOOTHING -------------
            snake = interpolate_points(snake, num_points)
            snake = remove_background_points(snake, b)
            snake = apply_gaussian_filter(snake, sigma=args.gauss_sigma)
            
            snake = snake[0:-1:10]

            # ------------- CALCULATE XY TORTUOSITY -------------
            L = np.sum(np.linalg.norm(np.diff(snake, axis=0), axis=-1))
            L_cumu = np.cumsum(np.linalg.norm(np.diff(snake, axis=0), axis=-1))
            L_0 = np.linalg.norm(snake[-1] - snake[0])
            xy_tortuosity = round(L / L_0, 4)
            # --------------------------------------------------------

            # calculate two boundaries. One at the 10th percentile of the snake and one at the 90th percentile.
            beginning = None
            end = (500,500)

            for i in range(len(snake)):
                if beginning is None and L_cumu[i] > L * 0.1:
                    beginning = snake[i]
                if L_cumu[i] > L * 0.9:
                    end = snake[i]
                    break
            vect = np.array([np.sin(-1 * rp.orientation), np.cos(rp.orientation)])
            bound1 = np.array([beginning - vect * 100, beginning + vect * 100])
            bound2 = np.array([end - vect * 100, end + vect * 100])

            print(f'bound1: {bound1}')

            # filter contours
            filtered_contours = filter_contours(contours, (beginning, end), rp.orientation, snake)

            # ------------- CALCULATE GRADE -------------
            b_gap_count = len(filtered_contours) - 1
            pred_grade = calculate_grade(xy_tortuosity, b_gap_count)
            print(f'xy_tortuosity: {xy_tortuosity}')
            # --------------------------------------------------------


            # create plot
            plt.imshow(image)
            fmts = ['c', '#ffa530', 'y', 'k', 'm', '#918211','#e875fa', '#b5d9ff', '#ffb5bb', '#3e754d', '#0014c7']
            print(f'# of filtered contours: {len(filtered_contours)}')
            for i, c in enumerate(filtered_contours):
                # plt.plot(c.T[1], c.T[0], f'-.',
                #          c=f'{fmts[i]}', lw=1, label=f'contour {i}')

                plt.fill(c.T[1], c.T[0], facecolor=f'{fmts[i]}', alpha=0.6, label=f'contour {i + 1}')
                

            # plt.scatter(snake[:, 1], snake[:, 0], marker='x', label='L', c='r', s=1)
            plt.plot(snake[:, 1], snake[:, 0], label='L', c='r', lw=1)
            plt.plot([snake[0,1], snake[-1,1]], [snake[0,0], snake[-1,0]], '--b', lw=1, label='L_0')
            plt.plot(bound1[:, 1], bound1[:, 0], 'g', lw=1, label='10th percentile of L')
            plt.plot(bound2[:, 1], bound2[:, 0], 'g', lw=1, label='90th percentile of L')
            plt.legend(fontsize=6)
            # plt.title(
            #     f'Filename: {fname}\nTortuosity (L/L_0): {xy_tortuosity}\nNumber of skips: {b_gap_count}')
            
            title = (
                f'Filename: {os.path.basename(fn)}'
                + (f' | Core Index: {core_mask_i}' if args.multiple_sections else '')
                + f'\nTortuosity (L/L_0): {xy_tortuosity}\nNumber of skips: {b_gap_count}\nPredicted Grade: Tier {pred_grade}\n'
            )
            plt.title(title)

            # create new slide
            blank_slide_layout = ppt.slide_layouts[6]
            slide = ppt.slides.add_slide(blank_slide_layout)


            # find largest index of greater tortuosity
            t_ind = len(tort_values) - 1
            new_loc_ind = -1

            if args.sort:
                while t_ind >= 0 and xy_tortuosity > tort_values[t_ind]:
                    new_loc_ind = t_ind
                    t_ind -= 1


            
            # move slide into sorted position if not already
            if new_loc_ind == -1:
                tort_values.append(xy_tortuosity)
                sorted_fns.append(fn)
                sorted_core_mask_indices.append(core_mask_i)
                sorted_skips.append(b_gap_count)
                sorted_pred_grades.append(pred_grade)
            else:
                tort_values.insert(new_loc_ind, xy_tortuosity)
                sorted_fns.insert(new_loc_ind, fn)
                sorted_core_mask_indices.insert(new_loc_ind, core_mask_i)
                sorted_skips.insert(new_loc_ind, b_gap_count)
                sorted_pred_grades.insert(new_loc_ind, pred_grade)
                move_slide(ppt.slides, slide, new_loc_ind)

            with BytesIO() as img_buf:
                # add predicted tortuosity to slide
                plt.savefig(img_buf, format='png', bbox_inches='tight', dpi=600.0)
                with Image.open(img_buf) as figure:
                    width, height = figure.size

                addimagetoslide(slide, 
                                img_buf, 
                                top=Inches(0), 
                                left=Inches(0), 
                                width=ppt.slide_width / 2, 
                                height=ppt.slide_width * (height/width) / 2)
                plt.close()

            # with BytesIO() as img_buf1:
            #     # add plot of curvature space to slide
            #     sigmas = np.linspace(0.1, 10, 50)
            #     res = curves.Curve(snake).scale_space(sigmas, features=['curvature'])

            #     fig = plt.figure(figsize=(7,3))
            #     ax = fig.add_subplot(111)
            #     pl = ax.imshow(np.log(res["curvature"]),cmap="jet")
            #     ax.invert_yaxis()
            #     ax.set_aspect(5)
            #     ax.grid()
            #     ax.set_xlabel("curve position index")
            #     ax.set_ylabel("sigma")
            #     cb = fig.colorbar(pl,shrink=0.7)
            #     cb.set_label("log curvature mu")
            #     plt.tight_layout()

            #     plt.savefig(img_buf1, format='png', bbox_inches='tight', dpi=600.0)
            #     with Image.open(img_buf1) as image:
            #         width, height = image.size

            #     addimagetoslide(slide, 
            #                     img_buf1, 
            #                     top=Inches(6), 
            #                     left=0, 
            #                     width=ppt.slide_width, 
            #                     height=ppt.slide_width * (height/width) / 2)
            #     plt.close()

            with BytesIO() as img_buf2:
                # add CNB thumbnail to slide
                osh = openslide.OpenSlide(fn)
                thumbnail = osh.get_thumbnail(size=(500, 500))
                
                plt.imshow(thumbnail)

                plt.savefig(img_buf2, format='png', bbox_inches='tight', dpi=600.0)
                with Image.open(img_buf2) as figure:
                    width, height = figure.size

                addimagetoslide(slide, 
                                img_buf2, 
                                top=Inches(0.5), 
                                left=ppt.slide_width/2, 
                                width=ppt.slide_width / 2, 
                                height=ppt.slide_width * (height/width) / 2)
                plt.close()

            if args.csv_save:
                # write sorted_fns and tort_values to a csv
                CORE_IDX_HEADER_INDEX = 1
                with open(args.csv_save, 'w', newline='') as csvfile:
                    writer = csv.writer(csvfile)
                    header_row = ['filename', 'tortuosity', '#skips', 'pred_grade']
                    if args.multiple_sections:
                        header_row.insert(CORE_IDX_HEADER_INDEX, 'core_index')
                    writer.writerow(header_row)
                    for i, fn in enumerate(sorted_fns):
                        row = [fn, tort_values[i], sorted_skips[i], sorted_pred_grades[i]]
                        if args.multiple_sections:
                            row.insert(CORE_IDX_HEADER_INDEX, sorted_core_mask_indices[i])
                        writer.writerow(row)


    ppt.save(args.ppt_save)

if __name__ == "__main__":
    main()